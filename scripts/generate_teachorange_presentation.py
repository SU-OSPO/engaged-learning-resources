#!/usr/bin/env python3
"""
Build a TeachOrange work presentation (PowerPoint).

Keynote is strict about OOXML. This script avoids:
- Custom slide sizes (use the template default)
- Manual slide XML (e.g. injected transitions)
- Blank layouts with extra shapes when a standard "Section Header" layout works

Usage (from repo root):
  source venv/bin/activate
  pip install python-pptx
  python scripts/generate_teachorange_presentation.py

Output: docs/TeachOrange_Work_Presentation.pptx

Add motion in Keynote: select slide → Animate → Add an Effect; or use PowerPoint.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

# Syracuse / TeachOrange palette (matches site)
NAVY = RGBColor(0x0D, 0x2D, 0x4C)
ORANGE = RGBColor(0xF7, 0x69, 0x00)
GRAY = RGBColor(0x55, 0x55, 0x55)

REPO_ROOT = Path(__file__).resolve().parent.parent
OUT_PATH = REPO_ROOT / "docs" / "TeachOrange_Work_Presentation.pptx"


def _add_bullets(text_frame, lines: list[str], *, level0_size: float = 20) -> None:
    text_frame.clear()
    first = True
    for line in lines:
        if first:
            p = text_frame.paragraphs[0]
            first = False
        else:
            p = text_frame.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(level0_size)
        p.font.color.rgb = GRAY
        p.space_after = Pt(6)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "TeachOrange"
    subtitle.text = (
        "Teaching activities library for higher-ed faculty\n"
        "Browse · Search · .edu access to materials\n\n"
        "Abhijnya Konanduru Gurumurthy\n"
        "Syracuse University · Maxwell School project"
    )
    for p in title.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = NAVY
            r.font.size = Pt(44)
            r.font.bold = True
    for p in subtitle.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = GRAY
    slide.notes_slide.notes_text_frame.text = (
        "Introduce the project goal: a shared library of teaching activities for faculty."
    )


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], notes: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    _add_bullets(body.text_frame, bullets)
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = NAVY
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_section_slide(prs: Presentation, heading: str, sub: str = "") -> None:
    """Standard Office 'Section Header' layout (Keynote-friendly)."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = heading
    sub_ph = slide.placeholders[1]
    sub_ph.text = sub or " "
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = NAVY
            r.font.bold = True
    for p in sub_ph.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(20)
            r.font.color.rgb = GRAY


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank you"
    sub = slide.placeholders[1]
    sub.text = (
        "Questions?\n\n"
        "Contact: Mark D. Brockway — mdbrockw@syr.edu\n"
        "TeachOrange · teachorange.com"
    )
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = ORANGE
            r.font.size = Pt(40)
    for p in sub.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(20)
            r.font.color.rgb = GRAY


def build() -> Path:
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    # Default template size (10 x 7.5 in — Keynote imports reliably)

    add_title_slide(prs)

    add_section_slide(prs, "Context", "Why TeachOrange exists")
    add_bullet_slide(
        prs,
        "Problem & audience",
        [
            "Faculty need quick access to peer-shared teaching activities.",
            "Materials (handouts, slides, prompts) should be easy to find and download.",
            "Access should stay faculty-focused (university email).",
        ],
        notes="Keep this slide short; emphasize higher-ed faculty as users.",
    )

    add_section_slide(prs, "What we shipped", "Product overview")
    add_bullet_slide(
        prs,
        "TeachOrange in one minute",
        [
            "Public catalog: browse and search activities by keyword, category, and tags.",
            "Detail pages: descriptions plus file previews/downloads for signed-in faculty.",
            "Self-service .edu registration, password reset via email (Mailjet).",
            "Django admin for professors to manage categories, tags, and uploads.",
        ],
    )

    add_section_slide(prs, "Experience highlights", "Features you can demo live")
    add_bullet_slide(
        prs,
        "Key interactions",
        [
            "Activity cards with thumbnail, summary, and tag chips.",
            "Click a tag on a card → filters the list to that tag.",
            "Preview supported file types in-browser after sign-in; download in one click.",
            "Responsive layout and consistent Syracuse-inspired branding (navy & orange).",
        ],
    )

    add_section_slide(prs, "Under the hood", "Engineering")
    add_bullet_slide(
        prs,
        "Technology stack",
        [
            "Backend: Django (Python) with PostgreSQL.",
            "Auth: faculty accounts tied to .edu addresses; protected media URLs.",
            "Email: transactional mail via Mailjet (e.g. password reset).",
            "Deployment path: production-ready for VPS (e.g. nginx + Gunicorn + TLS).",
        ],
        notes="Mention Django admin and Postgres if technical audience asks.",
    )

    add_section_slide(prs, "Team", "Who did what")
    add_bullet_slide(
        prs,
        "Credits",
        [
            "Faculty lead: Mark D. Brockway (Maxwell School, Syracuse University).",
            "Resource collection: Benjamin Katz & Sadie Rose Lehrfeld.",
            "Website & platform: Abhijnya Konanduru Gurumurthy (this presentation).",
        ],
    )

    add_section_slide(prs, "Status & next steps", "Roadmap")
    add_bullet_slide(
        prs,
        "Where we are",
        [
            "Core product: list, detail, auth, admin, tag filters, contact/team page.",
            "Next: production hosting, domain DNS, and ongoing content curation.",
            "Optional: more activities, analytics, and accessibility polish.",
        ],
    )

    add_bullet_slide(
        prs,
        "Presenter tip — animation",
        [
            "Keynote: sidebar → Animate → add an effect to the slide or to a text box.",
            "PowerPoint: Transitions / Animations tabs — try Fade or Appear on bullets.",
        ],
        notes="Optional slide; delete before presenting if you prefer.",
    )

    add_closing_slide(prs)

    prs.save(OUT_PATH)

    # Sanity check: file must round-trip in python-pptx
    Presentation(str(OUT_PATH))

    return OUT_PATH


def main() -> int:
    path = build()
    print(f"Wrote {path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
